"""
各格式文档：从字节提取纯文本/表格到编辑区，以及从编辑区内容生成 docx/xlsx/pdf 字节。
PPT 的解析与生成在 services.ppt_service / services.canvas_service。
"""

from __future__ import annotations

import io
import json
import re
from html.parser import HTMLParser
from typing import Any

from docx import Document
from docx.enum.text import WD_LINE_SPACING
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Mm, Pt, RGBColor
from openpyxl import Workbook
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation
from PyPDF2 import PdfReader
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfgen import canvas as rl_canvas


DOCX_FONT_CN = "Microsoft YaHei"
DOCX_FONT_MONO = "Consolas"

# Match frontend A4 paper: padding 25mm top/bottom, 20mm left/right
A4_MARGIN_TOP_MM = 25
A4_MARGIN_BOTTOM_MM = 25
A4_MARGIN_LEFT_MM = 20
A4_MARGIN_RIGHT_MM = 20

# Typography aligned with editor CSS (pt)
DOCX_BODY_SIZE_PT = 11
DOCX_H1_SIZE_PT = 22
DOCX_H2_SIZE_PT = 16
DOCX_H3_SIZE_PT = 14
DOCX_TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)
DOCX_BLOCKQUOTE_BORDER = "5B6B8C"
DOCX_FIRST_LINE_INDENT = Inches(0.35)
DOCX_BLOCKQUOTE_LEFT_INDENT = Inches(0.35)
# Hanging indent for manual bullets (bullet on line 1, text aligned beside it)
DOCX_LIST_LEFT_INDENT = Inches(0.35)
DOCX_LIST_HANGING_INDENT = Inches(-0.22)


def sanitize_text(s: str) -> str:
    if not s:
        return ""
    return re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", str(s))


def set_run_font(
    run: Any,
    name: str = DOCX_FONT_CN,
    *,
    size_pt: float | None = None,
    bold: bool | None = None,
    italic: bool | None = None,
    color: RGBColor | None = DOCX_TEXT_COLOR,
) -> None:
    """Set Latin + East Asia font so Chinese renders correctly in Word."""
    run.font.name = name
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.bold = bold
    if italic is not None:
        run.italic = italic
    if color is not None:
        run.font.color.rgb = color
    r_pr = run._element.get_or_add_rPr()
    r_fonts = r_pr.get_or_add_rFonts()
    for attr in ("w:ascii", "w:hAnsi", "w:eastAsia", "w:cs"):
        r_fonts.set(qn(attr), name)


def configure_a4_page(doc: Document) -> None:
    """A4 page size and margins matching the editor paper padding."""
    section = doc.sections[0]
    section.page_width = Mm(210)
    section.page_height = Mm(297)
    section.top_margin = Mm(A4_MARGIN_TOP_MM)
    section.bottom_margin = Mm(A4_MARGIN_BOTTOM_MM)
    section.left_margin = Mm(A4_MARGIN_LEFT_MM)
    section.right_margin = Mm(A4_MARGIN_RIGHT_MM)


def set_paragraph_shading(paragraph: Any, fill_hex: str) -> None:
    p_pr = paragraph._element.get_or_add_pPr()
    shading = OxmlElement("w:shd")
    shading.set(qn("w:val"), "clear")
    shading.set(qn("w:color"), "auto")
    shading.set(qn("w:fill"), fill_hex)
    p_pr.append(shading)


def set_paragraph_left_border(
    paragraph: Any,
    *,
    color: str = DOCX_BLOCKQUOTE_BORDER,
    size_eighths_pt: int = 24,
) -> None:
    """Left border ~4px, matching editor blockquote."""
    p_pr = paragraph._element.get_or_add_pPr()
    p_bdr = OxmlElement("w:pBdr")
    left = OxmlElement("w:left")
    left.set(qn("w:val"), "single")
    left.set(qn("w:sz"), str(size_eighths_pt))
    left.set(qn("w:space"), "8")
    left.set(qn("w:color"), color)
    p_bdr.append(left)
    p_pr.append(p_bdr)


def style_body_paragraph(paragraph: Any) -> None:
    paragraph.paragraph_format.space_after = Pt(6)
    paragraph.paragraph_format.line_spacing_rule = WD_LINE_SPACING.MULTIPLE
    paragraph.paragraph_format.line_spacing = 1.55


def style_normal_body_paragraph(paragraph: Any) -> None:
    """Body text with standard Chinese first-line indent (two characters)."""
    style_body_paragraph(paragraph)
    paragraph.paragraph_format.first_line_indent = DOCX_FIRST_LINE_INDENT


def style_blockquote_paragraph(paragraph: Any) -> None:
    """Left indent + accent border; no full-width background fill."""
    style_body_paragraph(paragraph)
    paragraph.paragraph_format.left_indent = DOCX_BLOCKQUOTE_LEFT_INDENT
    set_paragraph_left_border(
        paragraph,
        color=DOCX_BLOCKQUOTE_BORDER,
        size_eighths_pt=32,
    )
    for run in paragraph.runs:
        if run.italic is None:
            run.italic = True
        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def add_list_paragraph(
    doc: Document,
    text: str,
    *,
    ordered: bool = False,
    ol_index: int | None = None,
) -> Any:
    """
    One list line: bullet/number + text in a single paragraph.

    Avoids Word's built-in List Bullet style, which often splits the marker and
    CJK text onto separate lines when combined with first-line indent settings.
    """
    clean = sanitize_text(text)
    if not clean:
        return None
    if ordered and ol_index is not None:
        prefix = f"{ol_index}.\t"
    else:
        prefix = "•\t"
    paragraph = doc.add_paragraph()
    run = paragraph.add_run(prefix + clean)
    set_run_font(run, size_pt=DOCX_BODY_SIZE_PT)
    pf = paragraph.paragraph_format
    pf.space_after = Pt(4)
    pf.line_spacing_rule = WD_LINE_SPACING.MULTIPLE
    pf.line_spacing = 1.55
    pf.left_indent = DOCX_LIST_LEFT_INDENT
    pf.first_line_indent = DOCX_LIST_HANGING_INDENT
    return paragraph


def _flatten_li_inner_paragraphs(html: str) -> str:
    """
    Unwrap direct <p> children of <li> (Tiptap default).

    Uses a targeted pattern so nested <ul>/<li> trees are not broken.
    """
    pattern = re.compile(
        r"(<li\b[^>]*>)\s*<p\b[^>]*>(.*?)</p\s*>\s*(</li\s*>)",
        re.IGNORECASE | re.DOTALL,
    )
    prev = None
    while prev != html:
        prev = html
        html = pattern.sub(r"\1\2\3", html)
    return html


def add_body_heading_cn(doc: Document, text: str, level: int) -> Any:
    """Headings styled like the A4 editor (not Word default theme colors)."""
    sizes = {1: DOCX_H1_SIZE_PT, 2: DOCX_H2_SIZE_PT, 3: DOCX_H3_SIZE_PT}
    space_before = {1: Pt(14), 2: Pt(10), 3: Pt(8)}
    paragraph = doc.add_paragraph()
    run = paragraph.add_run(sanitize_text(text))
    set_run_font(run, size_pt=sizes.get(level, DOCX_BODY_SIZE_PT), bold=True)
    paragraph.paragraph_format.space_before = space_before.get(level, Pt(6))
    paragraph.paragraph_format.space_after = Pt(6 if level == 1 else 4)
    return paragraph


def set_paragraph_runs_font(
    paragraph: Any,
    name: str = DOCX_FONT_CN,
    *,
    size_pt: float | None = None,
) -> None:
    for run in paragraph.runs:
        set_run_font(run, name, size_pt=size_pt)


def apply_docx_chinese_defaults(doc: Document) -> None:
    """Configure document default style for CJK text."""
    normal = doc.styles["Normal"]
    normal.font.name = DOCX_FONT_CN
    normal.font.size = Pt(11)
    r_pr = normal._element.get_or_add_rPr()
    r_fonts = r_pr.get_or_add_rFonts()
    for attr in ("w:ascii", "w:hAnsi", "w:eastAsia", "w:cs"):
        r_fonts.set(qn(attr), DOCX_FONT_CN)


def add_heading_cn(doc: Document, text: str, level: int) -> Any:
    paragraph = doc.add_heading(sanitize_text(text), level=level)
    sizes = {0: 22, 1: 18, 2: 15, 3: 13}
    set_paragraph_runs_font(paragraph, size_pt=sizes.get(level, 12))
    return paragraph


def extract_pptx_text(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for si, slide in enumerate(prs.slides):
        lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text_frame.text.strip()
            if t:
                lines.append(t)
        if lines:
            parts.append(f"<!-- 幻灯片 {si + 1} -->\n" + "\n".join(lines))
    return "\n\n---\n\n".join(parts) if parts else "（未识别到文本，可在编辑区自行填写）"


def extract_docx_text(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    lines: list[str] = []
    for p in doc.paragraphs:
        t = sanitize_text(p.text)
        if not t:
            continue
        st = p.style.name if p.style else ""
        if "Heading 1" in st or st.startswith("Title"):
            lines.append(f"# {t}")
        elif "Heading 2" in st:
            lines.append(f"## {t}")
        elif "Heading 3" in st:
            lines.append(f"### {t}")
        elif p.style and "List" in p.style.name:
            lines.append(f"- {t}")
        else:
            lines.append(t)
    return "\n".join(lines) if lines else "（未识别到段落文本）"


def extract_xlsx_text(data: bytes) -> str:
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    ws = wb.active
    rows_out: list[str] = []
    for row in ws.iter_rows(values_only=True):
        cells = [("" if c is None else str(c)).replace("\t", " ").replace("\n", " ") for c in row]
        if any(cells):
            rows_out.append("\t".join(cells))
    wb.close()
    return "\n".join(rows_out) if rows_out else "列1\t列2\n示例\t数据"


def extract_pdf_text(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    chunks: list[str] = []
    for i, page in enumerate(reader.pages):
        try:
            t = page.extract_text() or ""
        except Exception:
            t = ""
        t = t.strip()
        if t:
            chunks.append(f"=== 第 {i + 1} 页 ===\n{t}")
    return "\n\n".join(chunks) if chunks else "（未能从 PDF 提取文本，可能是扫描件；请直接编辑）"


def _register_chinese_font() -> str:
    for fname in ("STSong-Light", "HeiseiKaku-W5", "HeiseiMin-W3"):
        try:
            pdfmetrics.registerFont(UnicodeCIDFont(fname))
            return fname
        except Exception:
            continue
    return "Helvetica"


def build_docx_bytes(topic: str, content: str) -> bytes:
    doc = Document()
    apply_docx_chinese_defaults(doc)
    add_heading_cn(doc, sanitize_text(topic.strip()) or "文档", 0)

    for raw in (content or "").splitlines():
        line = raw.rstrip()
        s = line.strip()
        if not s:
            continue
        if s.startswith("# ") and not s.startswith("##"):
            add_heading_cn(doc, s[2:], 1)
        elif s.startswith("## ") and not s.startswith("###"):
            add_heading_cn(doc, s[3:], 2)
        elif s.startswith("### "):
            add_heading_cn(doc, s[4:], 3)
        elif s.startswith("- ") or s.startswith("* "):
            try:
                p = doc.add_paragraph(sanitize_text(s[2:]), style="List Bullet")
            except KeyError:
                p = doc.add_paragraph()
                run = p.add_run("• " + sanitize_text(s[2:]))
                set_run_font(run)
            else:
                set_paragraph_runs_font(p)
            p.paragraph_format.space_after = Pt(4)
        elif s.startswith("<!--") and s.endswith("-->"):
            p = doc.add_paragraph()
            run = p.add_run(sanitize_text(s))
            set_run_font(run, italic=True)
        else:
            p = doc.add_paragraph()
            run = p.add_run(sanitize_text(s))
            set_run_font(run)
            p.paragraph_format.space_after = Pt(6)

    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()


class _TiptapHtmlToDocxParser(HTMLParser):
    """Parse Tiptap HTML into python-docx paragraphs (WYSIWYG-oriented)."""

    def __init__(self, document: Document) -> None:
        super().__init__(convert_charrefs=True)
        self.doc = document
        self._paragraph: Any | None = None
        self._heading_level: int | None = None
        self._in_pre = False
        self._pre_parts: list[str] = []
        self._bold = False
        self._italic = False
        self._in_blockquote = False
        self._in_li = False
        self._li_parts: list[str] = []
        self._list_stack: list[bool] = []
        self._ol_counter = 0
    def close(self) -> None:
        if self._in_li:
            self._flush_list_item()
        self._finish_paragraph()
        super().close()

    @staticmethod
    def _attrs_map(attrs: list[tuple[str, str | None]]) -> dict[str, str]:
        return {k: (v or "") for k, v in attrs if k}

    def _is_page_break_marker(self, tag: str, attrs: list[tuple[str, str | None]]) -> bool:
        amap = self._attrs_map(attrs)
        classes = amap.get("class", "").split()
        if tag == "hr":
            return "page-break" in classes
        if tag == "div":
            if "page-break" in classes:
                return True
            return amap.get("data-page-break", "").lower() in ("true", "1")
        return False

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        tag = tag.lower()
        if self._is_page_break_marker(tag, attrs):
            self._finish_paragraph()
            self.doc.add_page_break()
            return
        if tag in ("h1", "h2", "h3"):
            self._finish_paragraph()
            self._heading_level = int(tag[1])
        elif tag == "ul":
            self._list_stack.append(False)
        elif tag == "ol":
            self._list_stack.append(True)
            self._ol_counter = 0
        elif tag == "li":
            self._finish_paragraph()
            self._in_li = True
            self._li_parts = []
        elif tag == "p":
            if self._in_li:
                return
            self._finish_paragraph()
            self._paragraph = self.doc.add_paragraph()
        elif tag == "blockquote":
            self._finish_paragraph()
            self._in_blockquote = True
        elif tag == "pre":
            self._finish_paragraph()
            self._in_pre = True
            self._pre_parts = []
        elif tag in ("strong", "b"):
            self._bold = True
        elif tag in ("em", "i"):
            self._italic = True
        elif tag == "br":
            if self._in_li:
                self._li_parts.append(" ")
            elif self._paragraph is not None:
                self._paragraph.add_run().add_break()

    def handle_endtag(self, tag: str) -> None:
        tag = tag.lower()
        if tag == "hr":
            return
        if tag in ("h1", "h2", "h3"):
            self._heading_level = None
        elif tag in ("ul", "ol"):
            if self._list_stack:
                self._list_stack.pop()
        elif tag == "li":
            self._flush_list_item()
        elif tag == "p":
            if self._in_li:
                return
            if self._paragraph is not None:
                if self._in_blockquote:
                    style_blockquote_paragraph(self._paragraph)
                else:
                    style_normal_body_paragraph(self._paragraph)
            self._paragraph = None
        elif tag == "blockquote":
            self._in_blockquote = False
        elif tag == "pre":
            text = sanitize_text("".join(self._pre_parts))
            if text:
                paragraph = self.doc.add_paragraph()
                run = paragraph.add_run(text)
                set_run_font(run, DOCX_FONT_MONO, size_pt=10, color=DOCX_TEXT_COLOR)
                set_paragraph_shading(paragraph, "F6F8FA")
                paragraph.paragraph_format.space_after = Pt(6)
            self._in_pre = False
            self._pre_parts = []
        elif tag in ("strong", "b"):
            self._bold = False
        elif tag in ("em", "i"):
            self._italic = False

    def handle_data(self, data: str) -> None:
        if not data:
            return
        if self._in_pre:
            self._pre_parts.append(data)
            return
        if self._heading_level is not None:
            add_body_heading_cn(self.doc, data.strip(), self._heading_level)
            self._heading_level = None
            return
        if self._in_li:
            self._li_parts.append(data)
            return
        if self._paragraph is None:
            self._paragraph = self.doc.add_paragraph()
        run = self._paragraph.add_run(sanitize_text(data))
        use_italic = self._italic or self._in_blockquote
        set_run_font(
            run,
            bold=self._bold or None,
            italic=use_italic or None,
            color=RGBColor(0x55, 0x55, 0x55) if self._in_blockquote else DOCX_TEXT_COLOR,
        )

    def _flush_list_item(self) -> None:
        """Emit one list paragraph with full <li> text (ignores inner <p>)."""
        text = sanitize_text("".join(self._li_parts))
        ordered = self._list_stack[-1] if self._list_stack else False
        self._in_li = False
        self._li_parts = []
        if not text:
            return
        ol_index: int | None = None
        if ordered:
            self._ol_counter += 1
            ol_index = self._ol_counter
        add_list_paragraph(self.doc, text, ordered=ordered, ol_index=ol_index)

    def _finish_paragraph(self) -> None:
        if self._paragraph is not None:
            if self._in_blockquote:
                style_blockquote_paragraph(self._paragraph)
            elif not self._in_li:
                style_normal_body_paragraph(self._paragraph)
        self._paragraph = None
        self._heading_level = None


def append_html_to_docx(doc: Document, html: str, *, empty_label: str = "（空页）") -> None:
    """将一段 Tiptap HTML 追加到已有 Word 文档末尾。"""
    body = _flatten_li_inner_paragraphs((html or "").strip())
    if body:
        parser = _TiptapHtmlToDocxParser(doc)
        parser.feed(body)
        parser.close()
    else:
        p = doc.add_paragraph()
        run = p.add_run(empty_label)
        set_run_font(run)


def build_docx_from_html(html: str, topic: str) -> bytes:
    """Convert Tiptap editor HTML to an A4 Word document (WYSIWYG body only)."""
    del topic  # filename only; body matches editor canvas
    doc = Document()
    apply_docx_chinese_defaults(doc)
    configure_a4_page(doc)
    append_html_to_docx(doc, html, empty_label="（空文档）")

    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()


def build_docx_from_all_documents(
    pages: list[tuple[str, str]],
    *,
    book_title: str = "整本导出",
) -> bytes:
    """按顺序将多篇文档合并为一本 Word，篇与篇之间插入分页符。"""
    del book_title  # 文件名由路由层处理
    doc = Document()
    apply_docx_chinese_defaults(doc)
    configure_a4_page(doc)

    for index, (title, html) in enumerate(pages):
        if index > 0:
            doc.add_page_break()
        heading = sanitize_text((title or "").strip() or "未命名文档")
        h = doc.add_paragraph()
        run = h.add_run(heading)
        set_run_font(run, size_pt=DOCX_H1_SIZE_PT, bold=True)
        append_html_to_docx(doc, html)

    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()


def _parse_excel_grid(content: str) -> list[list[str]]:
    text = (content or "").strip()
    if not text:
        return [["列1", "列2"], ["", ""]]

    if text.startswith("{") and "columns" in text:
        try:
            obj = json.loads(text)
            cols = obj.get("columns") or []
            rows = obj.get("rows") or []
            if not isinstance(cols, list):
                cols = []
            if not isinstance(rows, list):
                rows = []
            grid: list[list[str]] = []
            if cols:
                grid.append([str(c) for c in cols])
            for r in rows:
                if isinstance(r, list):
                    grid.append([str(x) if x is not None else "" for x in r])
                else:
                    grid.append([str(r)])
            if not grid:
                return [["列1"], ["（无数据）"]]
            return grid
        except json.JSONDecodeError:
            pass

    grid: list[list[str]] = []
    for line in text.splitlines():
        if not line.strip():
            continue
        row = [sanitize_text(c) for c in line.split("\t")]
        if any(row):
            grid.append(row)
    if not grid:
        return [["列1", "列2"], ["（请在首行写表头，Tab 分列）", ""]]
    return grid


def build_xlsx_bytes(topic: str, content: str) -> bytes:
    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "Sheet1"
    grid = _parse_excel_grid(content)
    ws.append([sanitize_text(topic)])
    ws.append([])
    for row in grid:
        ws.append(row)
    for idx, col in enumerate(ws.columns, start=1):
        max_len = 0
        letter = get_column_letter(idx)
        for cell in col:
            try:
                max_len = max(max_len, len(str(cell.value or "")))
            except Exception:
                pass
        ws.column_dimensions[letter].width = min(max_len + 2, 48)

    bio = io.BytesIO()
    wb.save(bio)
    return bio.getvalue()


def build_pdf_bytes(topic: str, content: str) -> bytes:
    font = _register_chinese_font()
    bio = io.BytesIO()
    c = rl_canvas.Canvas(bio, pagesize=A4)
    w, h = A4
    margin = 2 * cm
    y = h - margin
    line_h = 14

    def new_page():
        nonlocal y
        c.showPage()
        y = h - margin

    def wrap_chunks(text: str, width: int = 36) -> list[str]:
        t = sanitize_text(text)
        if not t:
            return []
        return [t[i : i + width] for i in range(0, len(t), width)]

    def draw_line(text: str, size: int = 11, gap: int = 16):
        nonlocal y
        t = sanitize_text(text)
        if not t:
            return
        c.setFont(font, size)
        step = 18 if size >= 14 else 14
        for line in wrap_chunks(t, 36):
            if y < margin + line_h:
                new_page()
            c.drawString(margin, y, line)
            y -= step
        y -= gap

    draw_line(sanitize_text(topic.strip()) or "文档", size=16, gap=8)
    y -= 6

    for raw in (content or "").splitlines():
        s = raw.strip()
        if not s:
            y -= 6
            continue
        if s.startswith("# ") and not s.startswith("##"):
            draw_line(s[2:], size=14, gap=10)
        elif s.startswith("## "):
            draw_line(s[3:], size=13, gap=8)
        elif s.startswith("### "):
            draw_line(s[4:], size=12, gap=6)
        elif s.startswith("- ") or s.startswith("* "):
            draw_line("• " + s[2:], size=11, gap=4)
        elif s.startswith("==="):
            draw_line(s, size=12, gap=10)
        else:
            draw_line(s, size=11, gap=4)
        if y < margin:
            new_page()

    c.save()
    return bio.getvalue()


def sniff_format_from_filename(name: str) -> str | None:
    n = (name or "").lower()
    if n.endswith(".pptx"):
        return "ppt"
    if n.endswith(".docx"):
        return "word"
    if n.endswith(".xlsx") or n.endswith(".xlsm"):
        return "excel"
    if n.endswith(".pdf"):
        return "pdf"
    return None


def import_bytes_to_text(data: bytes, fmt: str) -> str:
    if fmt == "ppt":
        return extract_pptx_text(data)
    if fmt == "word":
        return extract_docx_text(data)
    if fmt == "excel":
        return extract_xlsx_text(data)
    if fmt == "pdf":
        return extract_pdf_text(data)
    raise ValueError(f"unsupported format {fmt}")


def export_bytes(fmt: str, topic: str, content: str) -> tuple[bytes, str, str]:
    t = sanitize_text(topic.strip()) or "document"
    if fmt == "word":
        return (
            build_docx_bytes(t, content),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            f"{t}_文档.docx",
        )
    if fmt == "excel":
        return (
            build_xlsx_bytes(t, content),
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            f"{t}_表格.xlsx",
        )
    if fmt == "pdf":
        return (
            build_pdf_bytes(t, content),
            "application/pdf",
            f"{t}_文档.pdf",
        )
    raise ValueError(f"unsupported export format {fmt}")
