"""Fabric 画布 ↔ PPTX / Word / Excel / PDF：导出与导入解析（所见即所得近似）。"""

from __future__ import annotations

import base64
import io
import re
from typing import Any, Iterator

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.oxml.ns import qn
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph as DocxParagraph
from docx.oxml import OxmlElement, parse_xml
from docx.shared import Inches, Length, Pt as DocPt, RGBColor as DocRGB
from openpyxl.styles import Alignment, Font as XLFont
from fastapi import HTTPException
from openpyxl import Workbook, load_workbook
from openpyxl.drawing.image import Image as XLImage
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Length, Pt
from PyPDF2 import PdfReader
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfgen import canvas as rl_canvas
from utils.file_utils import sanitize_text

# 与前端 Fabric 画布一致（16:9，96 DPI：960px = 10in）
CANVAS_W = 960.0
CANVAS_H = 540.0
SLIDE_W_EMU = 9144000
SLIDE_H_EMU = 5143500
SLIDE_W_IN = 10.0
SLIDE_H_IN = 5.625
SLIDE_W_PT = SLIDE_W_IN * 72.0
SLIDE_H_PT = SLIDE_H_IN * 72.0
PX_PER_IN = 96.0
# 与 static/js/canvas.js 导出常量一致
DEFAULT_LINE_HEIGHT = 1.35
DEFAULT_TEXT_PADDING_PX = 8.0
WORD_PAGE_MARGIN_PX = 40.0
PPT_STANDARD_MARGIN_PX = 48.0
# Word A4 画布（96 DPI 下 960px 宽 ≈ 10in，高度按 A4 比例）
WORD_CANVAS_W = 960.0
WORD_CANVAS_H = 1350.0
DOCX_DEFAULT_FONT_PT = 11.0
DOCX_TABLE_BORDER_PX = 1.0
# PPT 文本框默认内边距（未显式设置时）
PPT_TF_MARGIN_LR_EMU = 91440
PPT_TF_MARGIN_TB_EMU = 45720
TABLE_BORDER_EMU = 12700


def _hex_to_rgb(color: str | None) -> RGBColor | None:
    if not color or not isinstance(color, str):
        return None
    s = color.strip()
    if s.startswith("rgb"):
        m = re.search(r"rgba?\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)", s)
        if m:
            return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
        return None
    if not s.startswith("#"):
        return None
    h = s[1:]
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) != 6:
        return None
    try:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except ValueError:
        return None


def _rgb_to_hex(rgb) -> str:
    if rgb is None:
        return "#111111"
    try:
        r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
        return f"#{r:02x}{g:02x}{b:02x}"
    except Exception:
        return "#111111"


def _decode_image_src(src: str) -> bytes:
    if not src or not isinstance(src, str):
        raise ValueError("empty image src")
    if src.startswith("data:"):
        idx = src.find("base64,")
        if idx < 0:
            raise ValueError("invalid data URL")
        return base64.b64decode(src[idx + 7 :].strip())
    raise ValueError("only data:image base64 supported")


def _blank_slide_layout(prs: Presentation):
    for i in range(len(prs.slide_layouts) - 1, -1, -1):
        layout = prs.slide_layouts[i]
        try:
            if layout.name and "blank" in layout.name.lower():
                return layout
        except Exception:
            continue
    try:
        return prs.slide_layouts[6]
    except IndexError:
        return prs.slide_layouts[-1]


def _px_to_emu(val: float, canvas_px: float, slide_emu: int) -> int:
    """画布像素 → EMU，线性 1:1，无压缩。"""
    if canvas_px <= 0:
        return 0
    return int(round(float(val) * slide_emu / canvas_px))


def _emu_to_px(val: int | float, slide_emu: int, canvas_px: float) -> float:
    if not slide_emu:
        return 0.0
    return float(val) * canvas_px / float(slide_emu)


def _px_to_in(val: float, canvas_px: float, page_in: float) -> float:
    if canvas_px <= 0:
        return 0.0
    return float(val) * page_in / canvas_px


def _px_to_pt(val: float, canvas_px: float, page_pt: float) -> float:
    """画布字号(px) → Office 磅值，按页面宽度比例映射，与屏幕所见一致。"""
    if canvas_px <= 0:
        return float(val)
    return float(val) * page_pt / canvas_px


def _line_height_from_obj(obj: dict[str, Any]) -> float:
    try:
        lh = float(obj.get("lineHeight") or DEFAULT_LINE_HEIGHT)
    except (TypeError, ValueError):
        lh = DEFAULT_LINE_HEIGHT
    return max(1.0, min(lh, 3.0))


def _padding_px_from_obj(obj: dict[str, Any]) -> float:
    try:
        pad = float(obj.get("padding") if obj.get("padding") is not None else DEFAULT_TEXT_PADDING_PX)
    except (TypeError, ValueError):
        pad = DEFAULT_TEXT_PADDING_PX
    return max(0.0, min(pad, 48.0))


def _apply_ppt_text_frame_layout(
    tf,
    *,
    font_pt: float,
    fill: RGBColor | None,
    ff: str,
    line_height: float,
    padding_px: float,
    canvas_w: float,
    canvas_h: float,
) -> None:
    """文本框内边距、行距、段落间距，与 Fabric 预览对齐。"""
    pad_emu = _px_to_emu(padding_px, canvas_w, SLIDE_W_EMU)
    tf.margin_left = Emu(pad_emu)
    tf.margin_right = Emu(pad_emu)
    tf.margin_top = Emu(pad_emu)
    tf.margin_bottom = Emu(pad_emu)
    tf.word_wrap = True
    space_after_pt = max(2.0, font_pt * 0.08)
    for p in tf.paragraphs:
        p.line_spacing = line_height
        try:
            p.space_after = Pt(space_after_pt)
        except Exception:
            pass
        if fill is not None:
            p.font.color.rgb = fill
        p.font.size = Pt(font_pt)
        if ff:
            p.font.name = ff


def _mime_from_blob(blob: bytes) -> str:
    if len(blob) >= 8 and blob[:8] == b"\x89PNG\r\n\x1a\n":
        return "image/png"
    if len(blob) >= 2 and blob[:2] == b"\xff\xd8":
        return "image/jpeg"
    if len(blob) >= 6 and blob[:6] in (b"GIF87a", b"GIF89a"):
        return "image/gif"
    return "image/png"


def _iter_shapes_tree(shapes, gx: int = 0, gy: int = 0) -> Iterator[tuple[Any, int, int]]:
    for shp in shapes:
        ax = gx + int(shp.left or 0)
        ay = gy + int(shp.top or 0)
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes_tree(shp.shapes, ax, ay)
        else:
            yield shp, ax, ay


def _font_pt_from_run(font) -> float:
    try:
        if font is not None and font.size:
            return float(font.size.pt)
    except Exception:
        pass
    return 18.0


def _pt_to_canvas_font_px(pt: float, canvas_w: float = CANVAS_W) -> float:
    """PPT 磅值 → 画布 px（与导出 _px_to_pt 互逆）。"""
    return max(8.0, min(float(pt) * canvas_w / SLIDE_W_PT, 120.0))


def _paragraph_font_color_hex(p) -> str:
    for r in p.runs:
        if not r.text.strip():
            continue
        try:
            if r.font.color and r.font.color.rgb:
                return _rgb_to_hex(r.font.color.rgb)
        except Exception:
            continue
    try:
        if p.font.color and p.font.color.rgb:
            return _rgb_to_hex(p.font.color.rgb)
    except Exception:
        pass
    return "#111111"


def _paragraph_font_name(p) -> str:
    for r in p.runs:
        if r.text.strip() and r.font.name:
            return str(r.font.name)[:64]
    try:
        if p.font.name:
            return str(p.font.name)[:64]
    except Exception:
        pass
    return "Microsoft YaHei"


def _normalize_ppt_text(s: str) -> str:
    if not s:
        return ""
    return (
        s.replace("\x00", "")
        .replace("\x0b", "\n")
        .replace("\v", "\n")
        .replace("\r\n", "\n")
        .replace("\r", "\n")
    )


def _text_from_text_frame(tf) -> str:
    lines: list[str] = []
    for para in tf.paragraphs:
        lines.append(_normalize_ppt_text(para.text or ""))
    return "\n".join(lines).rstrip()


def _alignment_to_text_align(alignment: Any) -> str:
    if alignment is None:
        return "left"
    try:
        if alignment == PP_ALIGN.CENTER:
            return "center"
        if alignment == PP_ALIGN.RIGHT:
            return "right"
        if alignment in (PP_ALIGN.JUSTIFY, PP_ALIGN.JUSTIFY_LOW, PP_ALIGN.DISTRIBUTE, PP_ALIGN.THAI_DISTRIBUTE):
            return "justify"
    except Exception:
        pass
    return "left"


def _length_to_canvas_px(length: Length | None, slide_emu: int, canvas_px: float) -> float:
    if length is None:
        return 0.0
    try:
        return _emu_to_px(int(length), slide_emu, canvas_px)
    except Exception:
        return 0.0


def _line_spacing_to_fabric(ln_spc: Any, font_px: float) -> float:
    if ln_spc is None:
        return DEFAULT_LINE_HEIGHT
    try:
        if isinstance(ln_spc, Length):
            pt_val = float(ln_spc.pt)
            font_pt = font_px * SLIDE_W_PT / CANVAS_W if font_px > 0 else 12.0
            if font_pt > 0:
                return max(1.0, min(pt_val / font_pt, 3.0))
            return DEFAULT_LINE_HEIGHT
        return max(1.0, min(float(ln_spc), 3.0))
    except Exception:
        return DEFAULT_LINE_HEIGHT


def _margins_emu_from_text_frame(tf) -> tuple[int, int, int, int]:
    def one(attr: str, default: int) -> int:
        try:
            v = getattr(tf, attr, None)
            if v is None:
                return default
            return int(v)
        except Exception:
            return default

    return (
        one("margin_left", PPT_TF_MARGIN_LR_EMU),
        one("margin_right", PPT_TF_MARGIN_LR_EMU),
        one("margin_top", PPT_TF_MARGIN_TB_EMU),
        one("margin_bottom", PPT_TF_MARGIN_TB_EMU),
    )


def _margins_emu_from_cell(cell) -> tuple[int, int, int, int]:
    def one(attr: str, default: int) -> int:
        try:
            v = getattr(cell, attr, None)
            if v is None:
                return default
            return int(v)
        except Exception:
            return default

    return (
        one("margin_left", PPT_TF_MARGIN_LR_EMU),
        one("margin_right", PPT_TF_MARGIN_LR_EMU),
        one("margin_top", PPT_TF_MARGIN_TB_EMU),
        one("margin_bottom", PPT_TF_MARGIN_TB_EMU),
    )


def _padding_from_margins(
    ml: float, mr: float, mt: float, mb: float
) -> dict[str, float]:
    vals = [ml, mr, mt, mb]
    uniform = sum(vals) / 4.0
    return {
        "padding": round(uniform, 2),
        "paddingLeft": round(ml, 2),
        "paddingRight": round(mr, 2),
        "paddingTop": round(mt, 2),
        "paddingBottom": round(mb, 2),
    }


def _dominant_paragraph_style(tf) -> dict[str, Any]:
    chosen = None
    for para in tf.paragraphs:
        if (para.text or "").strip():
            chosen = para
            break
    if chosen is None and tf.paragraphs:
        chosen = tf.paragraphs[0]
    if chosen is None:
        return {
            "fontSize": _pt_to_canvas_font_px(18.0),
            "fill": "#111111",
            "fontFamily": "Microsoft YaHei",
            "textAlign": "left",
            "lineHeight": DEFAULT_LINE_HEIGHT,
        }
    r0 = next((r for r in chosen.runs if (r.text or "").strip()), chosen.runs[0] if chosen.runs else None)
    fs_pt = _font_pt_from_run(r0.font if r0 else chosen.font)
    font_px = _pt_to_canvas_font_px(fs_pt)
    return {
        "fontSize": font_px,
        "fill": _paragraph_font_color_hex(chosen),
        "fontFamily": _paragraph_font_name(chosen),
        "textAlign": _alignment_to_text_align(chosen.alignment),
        "lineHeight": _line_spacing_to_fabric(chosen.line_spacing, font_px),
    }


def _text_frame_to_canvas_object(
    tf,
    *,
    left_px: float,
    top_px: float,
    width_px: float,
    height_px: float,
    angle: float,
    sw: int,
    sh: int,
    margins_emu: tuple[int, int, int, int] | None = None,
) -> dict[str, Any] | None:
    text = _text_from_text_frame(tf)
    if not text:
        return None
    style = _dominant_paragraph_style(tf)
    ml_e, mr_e, mt_e, mb_e = margins_emu or _margins_emu_from_text_frame(tf)
    ml = _emu_to_px(ml_e, sw, CANVAS_W)
    mr = _emu_to_px(mr_e, sw, CANVAS_W)
    mt = _emu_to_px(mt_e, sh, CANVAS_H)
    mb = _emu_to_px(mb_e, sh, CANVAS_H)
    pad = _padding_from_margins(ml, mr, mt, mb)
    return {
        "type": "textbox",
        "left": left_px,
        "top": top_px,
        "width": max(8.0, width_px),
        "height": max(12.0, height_px),
        "text": sanitize_text(text)[:8000],
        "fontSize": style["fontSize"],
        "fill": style["fill"],
        "fontFamily": style["fontFamily"],
        "textAlign": style["textAlign"],
        "angle": angle,
        "lineHeight": style["lineHeight"],
        "fixedHeight": True,
        "splitByGrapheme": True,
        **pad,
    }


def _fill_to_hex(fill_obj) -> str | None:
    try:
        if fill_obj.type == MSO_FILL.SOLID:
            rgb = fill_obj.fore_color.rgb
            if rgb is not None:
                return _rgb_to_hex(rgb)
    except Exception:
        pass
    return None


def _slide_background_hex(slide) -> str:
    try:
        fill = slide.background.fill
        hx = _fill_to_hex(fill)
        if hx:
            return hx
    except Exception:
        pass
    return "#ffffff"


def _cell_vertical_offset_px(cell, inner_h_px: float, text_h_approx: float) -> float:
    try:
        anchor = cell.vertical_anchor
        if anchor == MSO_ANCHOR.MIDDLE:
            return max(0.0, (inner_h_px - text_h_approx) / 2.0)
        if anchor == MSO_ANCHOR.BOTTOM:
            return max(0.0, inner_h_px - text_h_approx)
    except Exception:
        pass
    return 0.0


def _parse_table_shape(
    shape,
    ax: int,
    ay: int,
    sw: int,
    sh: int,
    rot: float,
) -> list[dict[str, Any]]:
    tbl = shape.table
    objects: list[dict[str, Any]] = []
    ncols = len(tbl.columns)
    nrows = len(tbl.rows)
    if ncols < 1 or nrows < 1:
        return objects

    col_widths: list[int] = []
    for i in range(ncols):
        try:
            col_widths.append(int(tbl.columns[i].width or 0))
        except Exception:
            col_widths.append(0)
    row_heights: list[int] = []
    for i in range(nrows):
        try:
            row_heights.append(int(tbl.rows[i].height or 0))
        except Exception:
            row_heights.append(0)

    total_w = int(shape.width or 0) or sum(col_widths) or 1
    total_h = int(shape.height or 0) or sum(row_heights) or 1
    if sum(col_widths) <= 0:
        col_widths = [total_w // ncols] * ncols
        col_widths[0] += total_w - sum(col_widths)
    if sum(row_heights) <= 0:
        row_heights = [total_h // nrows] * nrows
        row_heights[0] += total_h - sum(row_heights)

    col_lefts = [0]
    for cw in col_widths[:-1]:
        col_lefts.append(col_lefts[-1] + cw)
    row_tops = [0]
    for rh in row_heights[:-1]:
        row_tops.append(row_tops[-1] + rh)

    border_px = max(1.0, _emu_to_px(TABLE_BORDER_EMU, sw, CANVAS_W))
    border_color = "#64748b"

    for ri in range(nrows):
        for ci in range(ncols):
            try:
                cell = tbl.cell(ri, ci)
            except Exception:
                continue
            if cell.is_spanned:
                continue

            span_w = int(cell.span_width) if cell.is_merge_origin else 1
            span_h = int(cell.span_height) if cell.is_merge_origin else 1
            cell_w = sum(col_widths[ci : ci + span_w])
            cell_h = sum(row_heights[ri : ri + span_h])
            cell_x = ax + col_lefts[ci]
            cell_y = ay + row_tops[ri]

            left_px = _emu_to_px(cell_x, sw, CANVAS_W)
            top_px = _emu_to_px(cell_y, sh, CANVAS_H)
            width_px = max(4.0, _emu_to_px(cell_w, sw, CANVAS_W))
            height_px = max(4.0, _emu_to_px(cell_h, sh, CANVAS_H))

            bg = _fill_to_hex(cell.fill)
            objects.append(
                {
                    "type": "rect",
                    "left": left_px,
                    "top": top_px,
                    "width": width_px,
                    "height": height_px,
                    "fill": bg or "#ffffff",
                    "stroke": border_color,
                    "strokeWidth": border_px,
                    "angle": rot,
                    "excludeFromExport": True,
                }
            )

            tf = cell.text_frame
            txt = _text_from_text_frame(tf)
            if not txt:
                continue

            margins = _margins_emu_from_cell(cell)
            style = _dominant_paragraph_style(tf)
            ml = _emu_to_px(margins[0], sw, CANVAS_W)
            mr = _emu_to_px(margins[1], sw, CANVAS_W)
            mt = _emu_to_px(margins[2], sh, CANVAS_H)
            mb = _emu_to_px(margins[3], sh, CANVAS_H)
            inner_h = max(8.0, height_px - mt - mb)
            line_h_px = style["fontSize"] * style["lineHeight"]
            line_count = max(1, txt.count("\n") + 1)
            v_off = _cell_vertical_offset_px(cell, inner_h, line_count * line_h_px)
            pad = _padding_from_margins(ml, mr, mt, mb)

            objects.append(
                {
                    "type": "textbox",
                    "left": left_px,
                    "top": top_px + mt + v_off,
                    "width": max(8.0, width_px),
                    "height": max(12.0, height_px - mt - mb),
                    "text": sanitize_text(txt)[:4000],
                    "fontSize": style["fontSize"],
                    "fill": style["fill"],
                    "fontFamily": style["fontFamily"],
                    "textAlign": style["textAlign"],
                    "angle": rot,
                    "lineHeight": style["lineHeight"],
                    "fixedHeight": True,
                    "splitByGrapheme": True,
                    **pad,
                }
            )

    return objects


def parse_pptx_to_canvas_payload(data: bytes) -> dict[str, Any]:
    """解析 PPTX：文本框/表格/图片按 EMU 线性映射到 960×540，还原边距、行高与对齐。"""
    prs = Presentation(io.BytesIO(data))
    sw = int(prs.slide_width or SLIDE_W_EMU)
    sh = int(prs.slide_height or SLIDE_H_EMU)
    slides_out: list[dict[str, Any]] = []

    for slide in prs.slides:
        objects: list[dict[str, Any]] = []
        for shape, ax, ay in _iter_shapes_tree(slide.shapes):
            try:
                rot = float(getattr(shape, "rotation", 0) or 0)
            except Exception:
                rot = 0.0

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                except Exception:
                    continue
                if not blob:
                    continue
                mime = _mime_from_blob(blob)
                src = "data:" + mime + ";base64," + base64.b64encode(blob).decode("ascii")
                w_emu = int(shape.width or 0)
                h_emu = int(shape.height or 0)
                objects.append(
                    {
                        "type": "image",
                        "left": _emu_to_px(ax, sw, CANVAS_W),
                        "top": _emu_to_px(ay, sh, CANVAS_H),
                        "width": max(8.0, _emu_to_px(w_emu, sw, CANVAS_W)),
                        "height": max(8.0, _emu_to_px(h_emu, sh, CANVAS_H)),
                        "angle": rot,
                        "src": src,
                    }
                )
                continue

            if getattr(shape, "has_table", False):
                objects.extend(_parse_table_shape(shape, ax, ay, sw, sh, rot))
                continue

            if shape.has_text_frame:
                tf = shape.text_frame
                w_emu = int(shape.width or 0)
                h_emu = int(shape.height or 0)
                obj = _text_frame_to_canvas_object(
                    tf,
                    left_px=_emu_to_px(ax, sw, CANVAS_W),
                    top_px=_emu_to_px(ay, sh, CANVAS_H),
                    width_px=max(40.0, _emu_to_px(w_emu, sw, CANVAS_W)),
                    height_px=max(24.0, _emu_to_px(h_emu, sh, CANVAS_H)),
                    angle=rot,
                    sw=sw,
                    sh=sh,
                )
                if obj:
                    objects.append(obj)

        slides_out.append({"bg": _slide_background_hex(slide), "objects": objects})

    if not slides_out:
        slides_out.append({"bg": "#ffffff", "objects": []})
    return {"canvasW": CANVAS_W, "canvasH": CANVAS_H, "slides": slides_out}


def _docx_len_to_px(length: Length | int | None, page_emu: int, canvas_px: float) -> float:
    if length is None or page_emu <= 0:
        return 0.0
    try:
        return float(int(length)) * canvas_px / float(page_emu)
    except Exception:
        return 0.0


def _estimate_wrapped_line_count(text: str, width_px: float, font_px: float) -> int:
    """按框宽估算 Word 段落自动换行后的行数（避免高度过小导致画布文字叠压）。"""
    if not text or not text.strip():
        return 1
    if width_px <= 0 or font_px <= 0:
        return max(1, text.count("\n") + 1)
    char_w = max(font_px * 0.52, 6.0)
    chars_per_line = max(1, int(width_px / char_w))
    total = 0
    for segment in text.replace("\r\n", "\n").split("\n"):
        n = len(segment)
        if n == 0:
            total += 1
        else:
            total += max(1, (n + chars_per_line - 1) // chars_per_line)
    return max(1, total)


def _docx_paragraph_block_height(
    text: str,
    width_px: float,
    font_px: float,
    line_height: float,
    *,
    space_after: float = 0.0,
    padding: float = DEFAULT_TEXT_PADDING_PX,
) -> float:
    lines = _estimate_wrapped_line_count(text, width_px, font_px)
    body = font_px * line_height * lines
    return body + padding * 2 + max(0.0, space_after) + 4.0


def _docx_pt_to_font_px(pt: float, page_w_emu: int, canvas_w: float = WORD_CANVAS_W) -> float:
    if page_w_emu <= 0:
        return max(8.0, min(float(pt), 120.0))
    page_w_pt = float(page_w_emu) / 12700.0
    if page_w_pt <= 0:
        return max(8.0, min(float(pt) * (canvas_w / 595.0), 120.0))
    return max(8.0, min(float(pt) * canvas_w / page_w_pt, 120.0))


def _iter_docx_blocks(document: Document) -> Iterator[Any]:
    """按 Word 文档顺序遍历段落与表格。"""
    parent_elm = document.element.body
    for child in parent_elm.iterchildren():
        if child.tag == qn("w:p"):
            yield DocxParagraph(child, document)
        elif child.tag == qn("w:tbl"):
            yield DocxTable(child, document)


def _docx_alignment_to_text_align(alignment: Any) -> str:
    if alignment is None:
        return "left"
    try:
        if alignment == WD_ALIGN_PARAGRAPH.CENTER:
            return "center"
        if alignment == WD_ALIGN_PARAGRAPH.RIGHT:
            return "right"
        if alignment in (WD_ALIGN_PARAGRAPH.JUSTIFY, WD_ALIGN_PARAGRAPH.DISTRIBUTE):
            return "justify"
    except Exception:
        pass
    return "left"


def _docx_line_height(pf, font_px: float) -> float:
    sp = pf.line_spacing
    rule = pf.line_spacing_rule
    if sp is None and rule is None:
        return DEFAULT_LINE_HEIGHT
    try:
        if rule == WD_LINE_SPACING.SINGLE:
            return 1.0
        if rule == WD_LINE_SPACING.ONE_POINT_FIVE:
            return 1.5
        if rule == WD_LINE_SPACING.DOUBLE:
            return 2.0
        if rule == WD_LINE_SPACING.MULTIPLE and sp is not None:
            return max(1.0, min(float(sp), 3.0))
        if rule == WD_LINE_SPACING.EXACTLY and sp is not None:
            pt_val = float(sp.pt) if hasattr(sp, "pt") else float(sp)
            if font_px > 0:
                return max(1.0, min(pt_val / (font_px * 12700.0 / WORD_CANVAS_W), 3.0))
        if hasattr(sp, "pt"):
            pt_val = float(sp.pt)
            if font_px > 0:
                return max(1.0, min(pt_val / font_px, 3.0))
        return max(1.0, min(float(sp), 3.0))
    except Exception:
        return DEFAULT_LINE_HEIGHT


def _docx_run_style(run) -> dict[str, Any]:
    fs_pt = DOCX_DEFAULT_FONT_PT
    fill = "#111111"
    ff = "Microsoft YaHei"
    if run is None:
        return {"font_pt": fs_pt, "fill": fill, "fontFamily": ff}
    try:
        if run.font.size:
            fs_pt = float(run.font.size.pt)
    except Exception:
        pass
    try:
        if run.font.color and run.font.color.rgb:
            fill = _rgb_to_hex(run.font.color.rgb)
    except Exception:
        pass
    try:
        if run.font.name:
            ff = str(run.font.name)[:64]
    except Exception:
        pass
    return {"font_pt": fs_pt, "fill": fill, "fontFamily": ff}


def _docx_paragraph_style(para: DocxParagraph, page_w_emu: int) -> dict[str, Any]:
    chosen = None
    for run in para.runs:
        if (run.text or "").strip():
            chosen = run
            break
    base = _docx_run_style(chosen) if chosen else _docx_run_style(para.runs[0] if para.runs else None)
    st = (para.style.name or "").lower()
    fs_pt = base["font_pt"]
    if "title" in st or "标题" in st:
        fs_pt = max(fs_pt, 26.0)
    elif "heading 1" in st or "标题 1" in st:
        fs_pt = max(fs_pt, 22.0)
    elif "heading 2" in st or "标题 2" in st:
        fs_pt = max(fs_pt, 18.0)
    elif "heading 3" in st or "标题 3" in st:
        fs_pt = max(fs_pt, 14.0)
    font_px = _docx_pt_to_font_px(fs_pt, page_w_emu)
    pf = para.paragraph_format
    return {
        "fontSize": font_px,
        "fill": base["fill"],
        "fontFamily": base["fontFamily"],
        "textAlign": _docx_alignment_to_text_align(pf.alignment),
        "lineHeight": _docx_line_height(pf, font_px),
    }


def _docx_run_has_page_break(run) -> bool:
    for br in run._element.findall(qn("w:br")):
        br_type = br.get(qn("w:type"))
        if br_type is None or br_type == "page":
            return True
    return False


def _docx_images_from_run(run, part, page_w_emu: int, page_h_emu: int) -> list[dict[str, Any]]:
    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    WP_NS = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
    out: list[dict[str, Any]] = []
    for drawing in run._element.iter():
        tag = drawing.tag
        if not tag.endswith("}inline") and not tag.endswith("}anchor"):
            continue
        cx, cy = 1905000, 1428750
        extent = drawing.find(f"{{{WP_NS}}}extent")
        if extent is not None:
            try:
                cx = int(extent.get("cx", cx))
                cy = int(extent.get("cy", cy))
            except (TypeError, ValueError):
                pass
        pos_x_emu, pos_y_emu = 0, 0
        if tag.endswith("}anchor"):
            pos_h = drawing.find(f"{{{WP_NS}}}positionH")
            pos_v = drawing.find(f"{{{WP_NS}}}positionV")
            if pos_h is not None:
                off = pos_h.find(f"{{{WP_NS}}}posOffset")
                if off is not None and off.text:
                    pos_x_emu = int(off.text)
            if pos_v is not None:
                off = pos_v.find(f"{{{WP_NS}}}posOffset")
                if off is not None and off.text:
                    pos_y_emu = int(off.text)
        for blip in drawing.iter(f"{{{A_NS}}}blip"):
            embed = blip.get(f"{{{R_NS}}}embed")
            if not embed:
                continue
            try:
                rel = part.related_parts[embed]
                blob = rel.blob
            except Exception:
                continue
            if not blob:
                continue
            mime = _mime_from_blob(blob)
            src = "data:" + mime + ";base64," + base64.b64encode(blob).decode("ascii")
            out.append(
                {
                    "src": src,
                    "width": max(8.0, _docx_len_to_px(cx, page_w_emu, WORD_CANVAS_W)),
                    "height": max(8.0, _docx_len_to_px(cy, page_h_emu, WORD_CANVAS_H)),
                    "pos_x": _docx_len_to_px(pos_x_emu, page_w_emu, WORD_CANVAS_W),
                    "pos_y": _docx_len_to_px(pos_y_emu, page_h_emu, WORD_CANVAS_H),
                    "anchored": tag.endswith("}anchor"),
                }
            )
    return out


def _docx_cell_fill_hex(cell) -> str | None:
    try:
        tc_pr = cell._element.tcPr
        if tc_pr is None:
            return None
        shd = tc_pr.find(qn("w:shd"))
        if shd is None:
            return None
        fill = shd.get(qn("w:fill"))
        if fill and fill.lower() not in ("auto", "ffffff", "ffffffff"):
            h = fill.lstrip("#")
            if len(h) == 6:
                return f"#{h}"
            if len(h) == 8:
                return f"#{h[2:]}"
    except Exception:
        pass
    return None


class _WordPageBuilder:
    """Word 分页画布构建器（A4 比例）。"""

    def __init__(self, page_w_emu: int, page_h_emu: int) -> None:
        self.page_w_emu = max(page_w_emu, 1)
        self.page_h_emu = max(page_h_emu, 1)
        self.canvas_w = WORD_CANVAS_W
        self.canvas_h = WORD_CANVAS_H
        self.margin_left = 72.0
        self.margin_right = 72.0
        self.margin_top = 72.0
        self.margin_bottom = 72.0
        self.content_left = self.margin_left
        self.content_width = self.canvas_w - self.margin_left - self.margin_right
        self.y = self.margin_top
        self.objects: list[dict[str, Any]] = []
        self.slides: list[dict[str, Any]] = []

    def set_margins_from_section(self, section) -> None:
        self.margin_left = _docx_len_to_px(section.left_margin, self.page_w_emu, self.canvas_w)
        self.margin_right = _docx_len_to_px(section.right_margin, self.page_w_emu, self.canvas_w)
        self.margin_top = _docx_len_to_px(section.top_margin, self.page_h_emu, self.canvas_h)
        self.margin_bottom = _docx_len_to_px(section.bottom_margin, self.page_h_emu, self.canvas_h)
        self.content_left = self.margin_left
        self.content_width = max(40.0, self.canvas_w - self.margin_left - self.margin_right)
        self.y = self.margin_top

    @property
    def content_bottom(self) -> float:
        return self.canvas_h - self.margin_bottom

    def flush_slide(self) -> None:
        if self.objects or not self.slides:
            self.slides.append({"bg": "#ffffff", "objects": self.objects[:]})
        self.objects = []
        self.y = self.margin_top

    def ensure_space(self, needed: float) -> None:
        if needed <= 0:
            return
        if self.y + needed > self.content_bottom and self.objects:
            self.flush_slide()

    def add_object(self, obj: dict[str, Any]) -> None:
        self.objects.append(obj)

    def add_image(self, img: dict[str, Any], default_left: float | None = None) -> None:
        w = float(img.get("width") or 80)
        h = float(img.get("height") or 60)
        if img.get("anchored"):
            left = float(img.get("pos_x") or 0) + self.margin_left
            top = float(img.get("pos_y") or 0) + self.margin_top
        else:
            left = default_left if default_left is not None else self.content_left
            top = self.y
            self.y += h + 4.0
        self.ensure_space(h)
        self.add_object(
            {
                "type": "image",
                "left": left,
                "top": top,
                "width": w,
                "height": h,
                "angle": 0.0,
                "src": img["src"],
            }
        )
        if img.get("anchored"):
            self.y = max(self.y, top + h + 4.0)

    def add_textbox(self, obj: dict[str, Any], advance_y: bool = True) -> None:
        self.add_object(obj)
        if advance_y:
            gap = float(obj.get("paraGap") or 2.0)
            self.y = float(obj.get("top", self.y)) + float(obj.get("height", 24)) + gap

    def _make_textbox(
        self,
        text: str,
        *,
        style: dict[str, Any],
        left: float,
        width: float,
        height: float,
        space_after: float = 0.0,
    ) -> dict[str, Any]:
        return {
            "type": "textbox",
            "left": left,
            "top": self.y,
            "width": max(8.0, width),
            "height": max(12.0, height),
            "text": sanitize_text(text)[:8000],
            "fontSize": style["fontSize"],
            "fill": style["fill"],
            "fontFamily": style["fontFamily"],
            "textAlign": style["textAlign"],
            "angle": 0.0,
            "lineHeight": style["lineHeight"],
            "padding": DEFAULT_TEXT_PADDING_PX,
            "fixedHeight": False,
            "splitByGrapheme": True,
            "spaceAfter": space_after,
            "paraGap": 2.0,
        }

    def parse_paragraph(self, para: DocxParagraph, doc: Document) -> None:
        pf = para.paragraph_format
        if pf.page_break_before:
            self.flush_slide()

        left_indent = _docx_len_to_px(pf.left_indent, self.page_w_emu, self.canvas_w)
        right_indent = _docx_len_to_px(pf.right_indent, self.page_w_emu, self.canvas_w)
        first_indent = _docx_len_to_px(pf.first_line_indent, self.page_w_emu, self.canvas_w)
        space_before = _docx_len_to_px(pf.space_before, self.page_h_emu, self.canvas_h)
        space_after = _docx_len_to_px(pf.space_after, self.page_h_emu, self.canvas_h)

        self.y += space_before
        box_left = self.content_left + left_indent + max(0.0, first_indent)
        box_width = max(20.0, self.content_width - left_indent - right_indent)

        part = doc.part
        for run in para.runs:
            if _docx_run_has_page_break(run):
                self.flush_slide()
            for img in _docx_images_from_run(run, part, self.page_w_emu, self.page_h_emu):
                self.add_image(img, default_left=box_left)

        text = _normalize_ppt_text(para.text or "")
        style = _docx_paragraph_style(para, self.page_w_emu)
        line_h = float(style["lineHeight"])
        font_px = float(style["fontSize"])

        if text.strip():
            box_h = _docx_paragraph_block_height(
                text, box_width, font_px, line_h, space_after=0.0
            )
            self.ensure_space(box_h + space_after)
            tb = self._make_textbox(
                text,
                style=style,
                left=box_left,
                width=box_width,
                height=box_h,
                space_after=space_after,
            )
            self.add_textbox(tb)
            if space_after > 0:
                self.y += space_after
        else:
            blank_h = max(space_after, font_px * line_h * 0.6, 6.0)
            self.y += blank_h

        for _pb in para.rendered_page_breaks:
            self.flush_slide()

    def parse_table(self, table: DocxTable) -> None:
        ncols = len(table.columns)
        nrows = len(table.rows)
        if ncols < 1 or nrows < 1:
            return

        col_widths: list[float] = []
        for col in table.columns:
            cw = col.width
            col_widths.append(
                _docx_len_to_px(cw, self.page_w_emu, self.canvas_w)
                if cw
                else self.content_width / ncols
            )
        if sum(col_widths) <= 0:
            col_widths = [self.content_width / ncols] * ncols

        row_heights: list[float] = []
        for row in table.rows:
            rh = row.height
            row_heights.append(
                _docx_len_to_px(rh, self.page_h_emu, self.canvas_h) if rh else max(22.0, 18.0 * 1.35)
            )

        table_h = sum(row_heights)
        self.ensure_space(table_h)
        top_base = self.y
        col_lefts = [0.0]
        for cw in col_widths[:-1]:
            col_lefts.append(col_lefts[-1] + cw)

        seen_tc: set[int] = set()
        border_color = "#64748b"

        for ri in range(nrows):
            for ci in range(ncols):
                try:
                    cell = table.cell(ri, ci)
                except Exception:
                    continue
                tc_id = id(cell._tc)
                if tc_id in seen_tc:
                    continue
                seen_tc.add(tc_id)

                cell_w = col_widths[ci] if ci < len(col_widths) else col_widths[-1]
                cell_h = row_heights[ri] if ri < len(row_heights) else row_heights[-1]
                left_px = self.content_left + col_lefts[ci]
                top_px = top_base + sum(row_heights[:ri])

                bg = _docx_cell_fill_hex(cell)
                self.add_object(
                    {
                        "type": "rect",
                        "left": left_px,
                        "top": top_px,
                        "width": max(4.0, cell_w),
                        "height": max(4.0, cell_h),
                        "fill": bg or "#ffffff",
                        "stroke": border_color,
                        "strokeWidth": DOCX_TABLE_BORDER_PX,
                        "angle": 0.0,
                        "excludeFromExport": True,
                    }
                )

                txt = _normalize_ppt_text(cell.text or "")
                if not txt.strip():
                    continue

                style = {
                    "fontSize": _docx_pt_to_font_px(DOCX_DEFAULT_FONT_PT, self.page_w_emu),
                    "fill": "#111111",
                    "fontFamily": "Microsoft YaHei",
                    "textAlign": "left",
                    "lineHeight": DEFAULT_LINE_HEIGHT,
                }
                for p_elm in cell._element.findall(qn("w:p")):
                    cp = DocxParagraph(p_elm, cell)
                    if (cp.text or "").strip():
                        style = _docx_paragraph_style(cp, self.page_w_emu)
                        break

                ml = mr = mt = mb = DEFAULT_TEXT_PADDING_PX
                try:
                    ml = _docx_len_to_px(cell.margin_left, self.page_w_emu, self.canvas_w)
                    mr = _docx_len_to_px(cell.margin_right, self.page_w_emu, self.canvas_w)
                    mt = _docx_len_to_px(cell.margin_top, self.page_h_emu, self.canvas_h)
                    mb = _docx_len_to_px(cell.margin_bottom, self.page_h_emu, self.canvas_h)
                except AttributeError:
                    pass
                font_px = float(style["fontSize"])
                lh = float(style["lineHeight"])
                inner_w = max(8.0, cell_w - ml - mr)
                inner_h = min(
                    max(12.0, cell_h - mt - mb),
                    _docx_paragraph_block_height(txt, inner_w, font_px, lh, padding=DEFAULT_TEXT_PADDING_PX),
                )

                self.add_object(
                    {
                        "type": "textbox",
                        "left": left_px + ml,
                        "top": top_px + mt,
                        "width": inner_w,
                        "height": max(12.0, inner_h),
                        "text": sanitize_text(txt)[:4000],
                        "fontSize": style["fontSize"],
                        "fill": style["fill"],
                        "fontFamily": style["fontFamily"],
                        "textAlign": style["textAlign"],
                        "angle": 0.0,
                        "lineHeight": style["lineHeight"],
                        "padding": DEFAULT_TEXT_PADDING_PX,
                        "fixedHeight": False,
                        "splitByGrapheme": True,
                    }
                )

        self.y = top_base + table_h + 8.0


def parse_docx_to_canvas_payload(data: bytes) -> dict[str, Any]:
    """解析 DOCX：A4 比例多页画布，还原段落/表格/图片与版式。"""
    doc = Document(io.BytesIO(data))
    section = doc.sections[0]
    page_w_emu = int(section.page_width or 7772400)
    page_h_emu = int(section.page_height or 10058400)
    builder = _WordPageBuilder(page_w_emu, page_h_emu)
    builder.set_margins_from_section(section)

    for block in _iter_docx_blocks(doc):
        if isinstance(block, DocxParagraph):
            builder.parse_paragraph(block, doc)
        elif isinstance(block, DocxTable):
            builder.parse_table(block)

    builder.flush_slide()
    slides = builder.slides
    if not slides:
        slides = [{"bg": "#ffffff", "objects": []}]
    return {"canvasW": WORD_CANVAS_W, "canvasH": WORD_CANVAS_H, "slides": slides}


def parse_xlsx_to_canvas_payload(data: bytes) -> dict[str, Any]:
    wb = load_workbook(io.BytesIO(data), data_only=True)
    ws = wb.active
    cell_w = 86.0
    cell_h = 24.0
    objects: list[dict[str, Any]] = []
    for row in ws.iter_rows():
        for cell in row:
            val = cell.value
            if val is None:
                continue
            s = sanitize_text(str(val).strip())
            if not s:
                continue
            ri = int(cell.row) - 1
            ci = int(cell.column) - 1
            left = 16.0 + ci * cell_w
            top = 16.0 + ri * cell_h
            if left > CANVAS_W - 30 or top > CANVAS_H - 20:
                continue
            objects.append(
                {
                    "type": "textbox",
                    "left": left,
                    "top": top,
                    "width": max(50.0, cell_w - 4),
                    "height": cell_h,
                    "text": s[:500],
                    "fontSize": 13.0,
                    "fill": "#0f172a",
                    "fontFamily": "Microsoft YaHei",
                    "angle": 0.0,
                }
            )
    if not objects:
        objects.append(
            {
                "type": "textbox",
                "left": 40.0,
                "top": 40.0,
                "width": 400.0,
                "height": 36.0,
                "text": "（空表或无法解析单元格）",
                "fontSize": 16.0,
                "fill": "#64748b",
                "fontFamily": "Microsoft YaHei",
                "angle": 0.0,
            }
        )
    wb.close()
    return {"canvasW": CANVAS_W, "canvasH": CANVAS_H, "slides": [{"bg": "#ffffff", "objects": objects}]}


def parse_pdf_to_canvas_payload(data: bytes) -> dict[str, Any]:
    reader = PdfReader(io.BytesIO(data))
    slides_out: list[dict[str, Any]] = []
    for page in reader.pages:
        try:
            raw = page.extract_text() or ""
        except Exception:
            raw = ""
        lines = [sanitize_text(ln.strip()) for ln in raw.splitlines() if sanitize_text(ln.strip())]
        objects: list[dict[str, Any]] = []
        y = 28.0
        for ln in lines[:48]:
            objects.append(
                {
                    "type": "textbox",
                    "left": 36.0,
                    "top": y,
                    "width": 880.0,
                    "height": 22.0,
                    "text": ln[:500],
                    "fontSize": 13.0,
                    "fill": "#111111",
                    "fontFamily": "Microsoft YaHei",
                    "angle": 0.0,
                }
            )
            y += 22.0
            if y > CANVAS_H - 16:
                break
        slides_out.append({"bg": "#ffffff", "objects": objects})
    if not slides_out:
        slides_out.append({"bg": "#ffffff", "objects": []})
    return {"canvasW": CANVAS_W, "canvasH": CANVAS_H, "slides": slides_out}


def build_ppt_from_canvas_payload(payload: dict[str, Any]) -> bytes:
    topic = str(payload.get("topic") or "课堂汇报").strip() or "课堂汇报"
    canvas_w = float(payload.get("canvasW") or CANVAS_W)
    canvas_h = float(payload.get("canvasH") or CANVAS_H)
    if canvas_w < 100 or canvas_h < 100:
        raise HTTPException(status_code=422, detail="canvasW/canvasH 过小")
    slides_in = payload.get("slides")
    if not isinstance(slides_in, list) or len(slides_in) < 1:
        raise HTTPException(status_code=422, detail="slides 必须为非空数组")

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    blank = _blank_slide_layout(prs)

    for si, slide_data in enumerate(slides_in):
        if not isinstance(slide_data, dict):
            continue
        slide = prs.slides.add_slide(blank)
        bg = slide_data.get("bg")
        rgb = _hex_to_rgb(str(bg)) if bg else None
        if rgb is not None:
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(SLIDE_W_EMU), Emu(SLIDE_H_EMU)
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = rgb
            bg_shape.line.fill.background()

        objects = slide_data.get("objects")
        if not isinstance(objects, list):
            objects = []

        for obj in objects:
            if not isinstance(obj, dict):
                continue
            typ = str(obj.get("type") or "").lower()
            left = _px_to_emu(float(obj.get("left") or 0), canvas_w, SLIDE_W_EMU)
            top = _px_to_emu(float(obj.get("top") or 0), canvas_h, SLIDE_H_EMU)
            width = max(_px_to_emu(float(obj.get("width") or 1), canvas_w, SLIDE_W_EMU), 9525)
            height = max(_px_to_emu(float(obj.get("height") or 1), canvas_h, SLIDE_H_EMU), 9525)
            angle = float(obj.get("angle") or 0)

            if typ == "image":
                try:
                    raw = _decode_image_src(str(obj.get("src") or ""))
                except Exception as e:
                    raise HTTPException(status_code=422, detail=f"幻灯片 {si + 1} 图片解码失败: {e}") from e
                stream = io.BytesIO(raw)
                pic = slide.shapes.add_picture(stream, Emu(left), Emu(top), width=Emu(width), height=Emu(height))
                pic.rotation = angle
            elif typ in ("textbox", "i-text", "text"):
                box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
                box.rotation = angle
                tf = box.text_frame
                tf.clear()
                text = str(obj.get("text") or "")
                lines = text.replace("\r\n", "\n").split("\n") if text else [""]
                try:
                    fs = float(obj.get("fontSize") or 24)
                except (TypeError, ValueError):
                    fs = 24.0
                font_pt = max(6.0, min(_px_to_pt(fs, canvas_w, SLIDE_W_PT), 96.0))
                fill = _hex_to_rgb(str(obj.get("fill") or "#111111"))
                ff = str(obj.get("fontFamily") or "").strip()[:64]
                lh = _line_height_from_obj(obj)
                pad_px = _padding_px_from_obj(obj)
                first = True
                for line in lines:
                    if first:
                        p = tf.paragraphs[0]
                        first = False
                    else:
                        p = tf.add_paragraph()
                    p.text = sanitize_text(line)
                _apply_ppt_text_frame_layout(
                    tf,
                    font_pt=font_pt,
                    fill=fill,
                    ff=ff,
                    line_height=lh,
                    padding_px=pad_px,
                    canvas_w=canvas_w,
                    canvas_h=canvas_h,
                )
            else:
                continue

    if len(prs.slides) == 0:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Emu(457200), Emu(257175), Emu(8000000), Emu(500000))
        box.text_frame.text = topic

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


def _word_emu(px: float, canvas_px: float, page_in: float) -> int:
    return int(round(_px_to_in(px, canvas_px, page_in) * 914400))


def _word_set_canvas_page(section, canvas_w: float, canvas_h: float) -> None:
    """Word 页面尺寸与 16:9 画布一致，边距与 parse_docx 默认 40px 对齐。"""
    section.page_width = Inches(SLIDE_W_IN)
    section.page_height = Inches(SLIDE_H_IN)
    margin_in = _px_to_in(WORD_PAGE_MARGIN_PX, canvas_w, SLIDE_W_IN)
    section.left_margin = Inches(margin_in)
    section.right_margin = Inches(margin_in)
    section.top_margin = Inches(_px_to_in(WORD_PAGE_MARGIN_PX, canvas_h, SLIDE_H_IN))
    section.bottom_margin = Inches(_px_to_in(WORD_PAGE_MARGIN_PX, canvas_h, SLIDE_H_IN))


def _xml_esc(text: str) -> str:
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def _word_add_floating_textbox(
    doc: Document,
    obj: dict[str, Any],
    canvas_w: float,
    canvas_h: float,
) -> None:
    """在 Word 中按画布坐标绝对定位文本框（页面坐标 1:1）。"""
    left = float(obj.get("left") or 0)
    top = float(obj.get("top") or 0)
    width = max(float(obj.get("width") or 80), 40.0)
    height = max(float(obj.get("height") or 24), 20.0)
    txt = sanitize_text(str(obj.get("text") or ""))
    if not txt:
        return
    try:
        fs = float(obj.get("fontSize") or 18)
    except (TypeError, ValueError):
        fs = 18.0
    font_pt = max(8.0, min(_px_to_pt(fs, canvas_w, SLIDE_W_PT), 72.0))
    fill_hex = str(obj.get("fill") or "#111111").lstrip("#")
    if len(fill_hex) == 3:
        fill_hex = "".join(c * 2 for c in fill_hex)
    ff = _xml_esc(str(obj.get("fontFamily") or "Microsoft YaHei").strip()[:64])
    lh = _line_height_from_obj(obj)
    pad_px = _padding_px_from_obj(obj)
    line_twips = int(font_pt * lh * 20)

    cx = _word_emu(width, canvas_w, SLIDE_W_IN)
    cy = _word_emu(height, canvas_h, SLIDE_H_IN)
    pos_x = _word_emu(left, canvas_w, SLIDE_W_IN)
    pos_y = _word_emu(top, canvas_h, SLIDE_H_IN)
    pad_l = _word_emu(pad_px, canvas_w, SLIDE_W_IN)
    pad_t = _word_emu(pad_px, canvas_h, SLIDE_H_IN)
    doc_pr_id = len(doc.element.body) + 1
    sz_half = int(round(font_pt * 2))

    paras_xml = ""
    for line in txt.replace("\r\n", "\n").split("\n"):
        paras_xml += (
            f'<w:p><w:pPr><w:spacing w:line="{line_twips}" w:lineRule="auto"/>'
            f'<w:ind w:left="0" w:right="0"/></w:pPr>'
            f'<w:r><w:rPr><w:sz w:val="{sz_half}"/><w:color w:val="{fill_hex}"/>'
            f'<w:rFonts w:ascii="{ff}" w:hAnsi="{ff}" w:eastAsia="{ff}"/></w:rPr>'
            f'<w:t xml:space="preserve">{_xml_esc(line)}</w:t></w:r></w:p>'
        )

    drawing_xml = f"""
    <w:drawing
      xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
      xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
      xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
      xmlns:wps="http://schemas.microsoft.com/office/word/2010/wordprocessingShape">
      <wp:anchor distT="0" distB="0" distL="0" distR="0" simplePos="0"
        relativeHeight="251658240" behindDoc="0" locked="0" layoutInCell="1" allowOverlap="1">
        <wp:simplePos x="0" y="0"/>
        <wp:positionH relativeFrom="page"><wp:posOffset>{pos_x}</wp:posOffset></wp:positionH>
        <wp:positionV relativeFrom="page"><wp:posOffset>{pos_y}</wp:posOffset></wp:positionV>
        <wp:extent cx="{cx}" cy="{cy}"/>
        <wp:effectExtent l="0" t="0" r="0" b="0"/>
        <wp:wrapNone/>
        <wp:docPr id="{doc_pr_id}" name="TextBox {doc_pr_id}"/>
        <wp:cNvGraphicFramePr/>
        <a:graphic>
          <a:graphicData uri="http://schemas.microsoft.com/office/word/2010/wordprocessingShape">
            <wps:wsp>
              <wps:cNvSpPr txBox="1"/>
              <wps:spPr>
                <a:xfrm><a:off x="0" y="0"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>
                <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
                <a:noFill/><a:ln w="0"><a:noFill/></a:ln>
              </wps:spPr>
              <wps:txbx><w:txbxContent>{paras_xml}</w:txbxContent></wps:txbx>
              <wps:bodyPr anchor="t" lIns="{pad_l}" tIns="{pad_t}" rIns="{pad_l}" bIns="{pad_t}"/>
            </wps:wsp>
          </a:graphicData>
        </a:graphic>
      </wp:anchor>
    </w:drawing>
    """
    p = doc.add_paragraph()
    run = p.add_run()
    run._r.append(parse_xml(drawing_xml))


def _word_add_floating_picture(
    doc: Document,
    obj: dict[str, Any],
    raw: bytes,
    canvas_w: float,
    canvas_h: float,
) -> None:
    """按画布坐标放置图片（绝对定位段落 + 精确尺寸）。"""
    left = float(obj.get("left") or 0)
    top = float(obj.get("top") or 0)
    width = max(float(obj.get("width") or 80), 20.0)
    height = max(float(obj.get("height") or 60), 20.0)
    p = doc.add_paragraph()
    pf = p.paragraph_format
    pf.space_before = DocPt(_px_to_pt(top, canvas_h, SLIDE_H_PT))
    pf.left_indent = DocPt(_px_to_pt(left, canvas_w, SLIDE_W_PT))
    pf.space_after = DocPt(0)
    pf.line_spacing_rule = WD_LINE_SPACING.SINGLE
    run = p.add_run()
    run.add_picture(
        io.BytesIO(raw),
        width=Inches(_px_to_in(width, canvas_w, SLIDE_W_IN)),
        height=Inches(_px_to_in(height, canvas_h, SLIDE_H_IN)),
    )


def _collect_canvas_objects(payload: dict[str, Any]) -> list[tuple[int, float, float, dict[str, Any]]]:
    """(slide_index, sort_y, sort_x, obj)"""
    slides_in = payload.get("slides") or []
    out: list[tuple[int, float, float, dict[str, Any]]] = []
    if not isinstance(slides_in, list):
        return out
    for si, slide in enumerate(slides_in):
        if not isinstance(slide, dict):
            continue
        objs = slide.get("objects")
        if not isinstance(objs, list):
            continue
        slide_base = float(si) * float(payload.get("canvasH") or CANVAS_H) * 10.0
        for obj in objs:
            if not isinstance(obj, dict):
                continue
            top = float(obj.get("top") or 0)
            left = float(obj.get("left") or 0)
            out.append((si, slide_base + top, left, obj))
    out.sort(key=lambda x: (x[0], x[1], x[2]))
    return out


def build_word_from_canvas_payload(payload: dict[str, Any]) -> bytes:
    canvas_w = float(payload.get("canvasW") or CANVAS_W)
    canvas_h = float(payload.get("canvasH") or CANVAS_H)
    doc = Document()
    _word_set_canvas_page(doc.sections[0], canvas_w, canvas_h)

    slides_in = payload.get("slides") or [{"bg": "#ffffff", "objects": []}]
    if not isinstance(slides_in, list):
        slides_in = [{"bg": "#ffffff", "objects": []}]

    for si, slide_data in enumerate(slides_in):
        if si > 0:
            doc.add_page_break()
        if not isinstance(slide_data, dict):
            continue
        objects = slide_data.get("objects")
        if not isinstance(objects, list):
            continue
        ordered = sorted(
            objects,
            key=lambda o: (float(o.get("top") or 0), float(o.get("left") or 0))
            if isinstance(o, dict)
            else (0.0, 0.0),
        )
        for obj in ordered:
            if not isinstance(obj, dict):
                continue
            typ = str(obj.get("type") or "").lower()
            if typ == "image":
                try:
                    raw = _decode_image_src(str(obj.get("src") or ""))
                except Exception:
                    continue
                _word_add_floating_picture(doc, obj, raw, canvas_w, canvas_h)
            elif typ in ("textbox", "i-text", "text"):
                _word_add_floating_textbox(doc, obj, canvas_w, canvas_h)

    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()


def build_excel_from_canvas_payload(payload: dict[str, Any]) -> bytes:
    cw = float(payload.get("canvasW") or CANVAS_W)
    ch = float(payload.get("canvasH") or CANVAS_H)
    col_px = 48.0
    row_px = 18.0
    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "Canvas"
    ws.sheet_view.showGridLines = False

    max_col = max(1, int(cw / col_px) + 2)
    max_row = max(1, int(ch / row_px) + 2)
    for col in range(1, max_col + 1):
        ws.column_dimensions[get_column_letter(col)].width = max(6.0, col_px / 7.0)
    for row in range(1, max_row + 1):
        ws.row_dimensions[row].height = max(12.0, row_px * 0.75)

    slides_in = payload.get("slides") or [{"objects": []}]
    if not isinstance(slides_in, list):
        slides_in = [{"objects": []}]

    row_offset = 0
    for si, slide_data in enumerate(slides_in):
        if si > 0:
            row_offset += int(ch / row_px) + 2
        if not isinstance(slide_data, dict):
            continue
        objects = slide_data.get("objects")
        if not isinstance(objects, list):
            continue
        for obj in objects:
            if not isinstance(obj, dict):
                continue
            typ = str(obj.get("type") or "").lower()
            left = float(obj.get("left") or 0)
            top = float(obj.get("top") or 0) + row_offset * row_px
            ci = max(1, int(left / col_px) + 1)
            ri = max(1, int(top / row_px) + 1)

            if typ == "image":
                try:
                    raw = _decode_image_src(str(obj.get("src") or ""))
                except Exception:
                    continue
                img = XLImage(io.BytesIO(raw))
                img.width = int(max(20, min(float(obj.get("width") or 120), cw)))
                img.height = int(max(16, min(float(obj.get("height") or 80), ch)))
                ws.add_image(img, f"{get_column_letter(ci)}{ri}")
                continue

            if typ not in ("textbox", "i-text", "text"):
                continue
            txt = sanitize_text(str(obj.get("text") or ""))
            if not txt:
                continue
            try:
                fs = float(obj.get("fontSize") or 13)
            except (TypeError, ValueError):
                fs = 13.0
            font_pt = max(8.0, min(_px_to_pt(fs, cw, SLIDE_W_PT), 44.0))
            lh = _line_height_from_obj(obj)
            cell = ws.cell(row=ri, column=ci, value=txt)
            cell.font = XLFont(
                name=str(obj.get("fontFamily") or "Microsoft YaHei")[:64],
                size=font_pt,
            )
            cell.alignment = Alignment(
                wrap_text=True,
                vertical="top",
                horizontal="left",
                indent=max(0, int(_padding_px_from_obj(obj) / 4)),
            )
            span_h = max(1, int(float(obj.get("height") or row_px) / row_px))
            ws.row_dimensions[ri].height = max(
                ws.row_dimensions[ri].height or 12.0,
                row_px * span_h * lh * 0.75,
            )

    bio = io.BytesIO()
    wb.save(bio)
    return bio.getvalue()


def _pdf_register_font() -> str:
    for fname in ("STSong-Light", "HeiseiKaku-W5", "HeiseiMin-W3"):
        try:
            pdfmetrics.registerFont(UnicodeCIDFont(fname))
            return fname
        except Exception:
            continue
    return "Helvetica"


def build_pdf_from_canvas_payload(payload: dict[str, Any]) -> bytes:
    canvas_w = float(payload.get("canvasW") or CANVAS_W)
    canvas_h = float(payload.get("canvasH") or CANVAS_H)
    font = _pdf_register_font()
    bio = io.BytesIO()
    page_size = (canvas_w, canvas_h)
    c = rl_canvas.Canvas(bio, pagesize=page_size)
    pad = DEFAULT_TEXT_PADDING_PX

    slides_in = payload.get("slides") or []
    if not isinstance(slides_in, list) or not slides_in:
        slides_in = [{"bg": "#ffffff", "objects": []}]

    for si, slide in enumerate(slides_in):
        if not isinstance(slide, dict):
            continue
        if si > 0:
            c.showPage()
        bg = slide.get("bg")
        if bg and isinstance(bg, str):
            rgb = _hex_to_rgb(bg)
            if rgb is not None:
                c.setFillColorRGB(rgb[0] / 255.0, rgb[1] / 255.0, rgb[2] / 255.0)
                c.rect(0, 0, canvas_w, canvas_h, fill=1, stroke=0)
                c.setFillColorRGB(0.07, 0.09, 0.15)

        objs = slide.get("objects")
        if not isinstance(objs, list):
            continue
        for obj in objs:
            if not isinstance(obj, dict):
                continue
            typ = str(obj.get("type") or "").lower()
            left = float(obj.get("left") or 0) + pad
            top_canvas = float(obj.get("top") or 0)
            box_h = float(obj.get("height") or 24)

            if typ in ("textbox", "i-text", "text"):
                try:
                    fs = float(obj.get("fontSize") or 12)
                except (TypeError, ValueError):
                    fs = 12.0
                lh = _line_height_from_obj(obj)
                line_leading = fs * lh
                fill_rgb = _hex_to_rgb(str(obj.get("fill") or "#111111"))
                if fill_rgb is not None:
                    c.setFillColorRGB(fill_rgb[0] / 255.0, fill_rgb[1] / 255.0, fill_rgb[2] / 255.0)
                c.setFont(font, fs)
                txt = sanitize_text(str(obj.get("text") or ""))
                y_base = canvas_h - top_canvas - pad - fs
                for i, line in enumerate(txt.replace("\r\n", "\n").split("\n")):
                    c.drawString(left, y_base - i * line_leading, line[:500])
            elif typ == "image":
                try:
                    raw = _decode_image_src(str(obj.get("src") or ""))
                except Exception:
                    continue
                w = float(obj.get("width") or 100)
                h = float(obj.get("height") or 80)
                y_img = canvas_h - top_canvas - h
                try:
                    c.drawImage(io.BytesIO(raw), float(obj.get("left") or 0), y_img, width=w, height=h)
                except Exception:
                    continue
    c.save()
    return bio.getvalue()


def build_doc_canvas_file(fmt: str, payload: dict[str, Any]) -> tuple[bytes, str, str]:
    t = sanitize_text(str(payload.get("topic") or "导出").strip()) or "导出"
    if fmt == "word":
        return (
            build_word_from_canvas_payload(payload),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            f"{t}_画布.docx",
        )
    if fmt == "excel":
        return (
            build_excel_from_canvas_payload(payload),
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            f"{t}_画布.xlsx",
        )
    if fmt == "pdf":
        return (
            build_pdf_from_canvas_payload(payload),
            "application/pdf",
            f"{t}_画布.pdf",
        )
    raise ValueError("unsupported canvas doc format")
