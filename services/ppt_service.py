"""Markdown / 结构化大纲解析与 python-pptx 多页生成。"""

from __future__ import annotations

import io
import re
import urllib.parse
from typing import Any

from fastapi import HTTPException
from pptx import Presentation
from pptx.util import Emu, Pt

MAX_BODY_ENTRIES_PER_SLIDE = 8
DEFAULT_LINE_HEIGHT = 1.35
BODY_FONT_PT = 18
SUB_FONT_PT = 16
TEXT_FRAME_MARGIN_EMU = 45720  # ~0.05 inch，避免占位符内文字贴边


def sanitize_ooxml_text(s: str) -> str:
    if not s:
        return ""
    return re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", s)


def parse_custom_outline_to_data(topic: str, custom_outline: str) -> dict[str, Any]:
    text = (custom_outline or "").strip()
    if not text:
        raise HTTPException(status_code=422, detail="custom_outline 不能为空")

    subtitle = "课堂汇报"
    chapters: list[dict[str, Any]] = []
    in_toc = False
    current_ch: dict[str, Any] | None = None
    current_sec: dict[str, Any] | None = None

    def clean(s: str) -> str:
        return sanitize_ooxml_text((s or "").strip())

    def ensure_ch() -> dict[str, Any]:
        nonlocal current_ch, chapters
        if current_ch is None:
            current_ch = {"title": "正文", "sections": []}
            chapters.append(current_ch)
        return current_ch

    def ensure_sec() -> dict[str, Any]:
        nonlocal current_sec
        ch = ensure_ch()
        if current_sec is None:
            current_sec = {"heading": "要点", "points": []}
            ch["sections"].append(current_sec)
        return current_sec

    for raw in text.splitlines():
        stripped = raw.strip()

        m_sub = re.match(r"^\*{0,2}副标题\*{0,2}[：:]\s*(.+)$", stripped)
        if m_sub:
            subtitle = clean(m_sub.group(1)) or subtitle
            continue

        if stripped.startswith("#") and not stripped.startswith("##"):
            continue

        if stripped.startswith("## "):
            rest = clean(stripped[3:])
            if rest == "目录" or rest.startswith("目录"):
                in_toc = True
                current_ch = None
                current_sec = None
                continue
            in_toc = False
            current_sec = None
            if not rest:
                current_ch = None
                continue
            current_ch = {"title": rest, "sections": []}
            chapters.append(current_ch)
            continue

        if in_toc:
            continue

        if stripped.startswith("### "):
            ch = ensure_ch()
            current_sec = {"heading": clean(stripped[4:]) or "要点", "points": []}
            ch["sections"].append(current_sec)
            continue

        if stripped.startswith("- ") or stripped.startswith("* "):
            pt = clean(stripped[2:])
            if pt:
                ensure_sec()["points"].append(pt)
            continue

        if stripped and not stripped.startswith("#") and "**副标题**" not in stripped:
            pt = clean(stripped)
            if pt:
                ensure_sec()["points"].append(pt)

    if not chapters:
        pts = []
        for raw in text.splitlines():
            s = clean(raw)
            if not s or s.startswith("#"):
                continue
            pts.append(s)
        pts = pts[:60]
        if not pts:
            raise HTTPException(
                status_code=422,
                detail="无法解析大纲：请至少包含 ## 章节标题，或使用 ## / ### / - 列表格式书写。",
            )
        chapters = [{"title": "正文", "sections": [{"heading": "内容", "points": pts}]}]

    for ch in chapters:
        secs = ch.get("sections") or []
        if not secs:
            ch["sections"] = [
                {"heading": "要点", "points": ["（可在编辑区为本章补充 ### 小标题与 - 要点）"]}
            ]
        for sec in ch["sections"]:
            if not isinstance(sec, dict):
                continue
            sec["heading"] = clean(str(sec.get("heading", "要点")))
            pts = sec.get("points") or []
            sec["points"] = [clean(str(p)) for p in pts if clean(str(p))]
            if not sec["points"] and sec.get("heading"):
                sec["points"] = ["（本节暂无要点，可在编辑区添加 - 列表项）"]

    subtitle = clean(subtitle) or "课堂汇报"
    return {"subtitle": subtitle, "chapters": chapters}


def outline_json_to_markdown(topic: str, data: dict[str, Any]) -> str:
    lines: list[str] = [f"# {topic}", "", f"**副标题**：{data.get('subtitle', '课堂汇报')}", "", "## 目录"]
    for ch in data.get("chapters", []):
        if isinstance(ch, dict) and ch.get("title"):
            lines.append(f"- {ch['title']}")
    lines.append("")
    for ch in data.get("chapters", []):
        if not isinstance(ch, dict):
            continue
        title = ch.get("title", "")
        lines.append(f"## {title}")
        for sec in ch.get("sections", []) or []:
            if not isinstance(sec, dict):
                continue
            lines.append(f"### {sec.get('heading', '')}")
            for p in sec.get("points", []) or []:
                lines.append(f"- {p}")
            lines.append("")
    return "\n".join(lines).strip()


def _flatten_chapter_body(chapter: dict[str, Any]) -> list[tuple[int, str]]:
    rows: list[tuple[int, str]] = []
    for sec in chapter.get("sections", []) or []:
        if not isinstance(sec, dict):
            continue
        h = (sec.get("heading") or "").strip()
        if h:
            rows.append((0, h))
        for pt in sec.get("points", []) or []:
            t = str(pt).strip()
            if t:
                rows.append((1, t))
    return rows


def _chunk_rows(rows: list[tuple[int, str]], max_per_slide: int) -> list[list[tuple[int, str]]]:
    if not rows:
        return [[]]
    chunks: list[list[tuple[int, str]]] = []
    cur: list[tuple[int, str]] = []
    for row in rows:
        if cur and len(cur) >= max_per_slide:
            chunks.append(cur)
            cur = []
        cur.append(row)
    if cur:
        chunks.append(cur)
    return chunks


def _set_slide_title_content(
    slide,
    title: str,
    rows: list[tuple[int, str]],
) -> None:
    slide.shapes.title.text = sanitize_ooxml_text(title)
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    tf.margin_left = Emu(TEXT_FRAME_MARGIN_EMU)
    tf.margin_right = Emu(TEXT_FRAME_MARGIN_EMU)
    tf.margin_top = Emu(TEXT_FRAME_MARGIN_EMU)
    tf.margin_bottom = Emu(TEXT_FRAME_MARGIN_EMU)
    tf.word_wrap = True

    if not rows:
        p = tf.paragraphs[0]
        p.text = "（本章节暂无要点）"
        p.level = 0
        p.line_spacing = DEFAULT_LINE_HEIGHT
        p.space_after = Pt(6)
        return

    first = True
    for level, text in rows:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = sanitize_ooxml_text(text)
        p.level = level
        pt = BODY_FONT_PT if level == 0 else SUB_FONT_PT
        p.font.size = Pt(pt)
        p.line_spacing = DEFAULT_LINE_HEIGHT
        p.space_after = Pt(6 if level == 0 else 4)


def build_ppt_bytes(topic: str, data: dict[str, Any]) -> bytes:
    prs = Presentation()
    topic_safe = sanitize_ooxml_text(topic.strip()) or "课堂汇报"
    subtitle = sanitize_ooxml_text(str(data.get("subtitle") or "课堂汇报").strip()) or "课堂汇报"
    chapters = [c for c in (data.get("chapters") or []) if isinstance(c, dict) and (c.get("title") or "").strip()]

    slide0 = prs.slides.add_slide(prs.slide_layouts[0])
    slide0.shapes.title.text = topic_safe
    if len(slide0.placeholders) > 1:
        slide0.placeholders[1].text = subtitle

    toc = prs.slides.add_slide(prs.slide_layouts[1])
    toc.shapes.title.text = "目录"
    tf = toc.placeholders[1].text_frame
    tf.clear()
    tf.margin_left = Emu(TEXT_FRAME_MARGIN_EMU)
    tf.margin_right = Emu(TEXT_FRAME_MARGIN_EMU)
    tf.margin_top = Emu(TEXT_FRAME_MARGIN_EMU)
    tf.margin_bottom = Emu(TEXT_FRAME_MARGIN_EMU)
    for i, ch in enumerate(chapters):
        title_line = str(ch.get("title", "")).strip()
        if not title_line:
            continue
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = sanitize_ooxml_text(title_line)
        p.level = 0
        p.font.size = Pt(20)
        p.line_spacing = DEFAULT_LINE_HEIGHT
        p.space_after = Pt(8)

    for ch in chapters:
        title = str(ch.get("title", "")).strip()
        rows = _flatten_chapter_body(ch)
        chunks = _chunk_rows(rows, MAX_BODY_ENTRIES_PER_SLIDE)
        if not chunks:
            s = prs.slides.add_slide(prs.slide_layouts[1])
            _set_slide_title_content(s, title, [])
            continue
        for idx, chunk in enumerate(chunks):
            page_title = title if idx == 0 else f"{title}（续）"
            s = prs.slides.add_slide(prs.slide_layouts[1])
            _set_slide_title_content(s, page_title, chunk)

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


def attachment_headers(*, ascii_filename: str, utf8_filename: str) -> dict[str, str]:
    filename_star = urllib.parse.quote(utf8_filename, safe="")
    return {
        "Content-Disposition": (
            f'attachment; filename="{ascii_filename}"; filename*=UTF-8\'\'{filename_star}'
        ),
    }


def ascii_stem(topic: str) -> str:
    stem = re.sub(r"[^\x20-\x7E]", "_", (topic[:80] or "topic")).strip() or "topic"
    stem = re.sub(r'[<>:"/\\|?*\s]', "_", stem).strip("_") or "topic"
    return stem[:72]
